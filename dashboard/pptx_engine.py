# ===========================================================================
# PPTX ENGINE: PowerPoint export for storyline decks
# ===========================================================================
# Generates branded PPTX files from storyline configs using python-pptx.
# Each slide receives a chart image (rendered from matplotlib figures),
# a title bar, subtitle, and footer.

import os
import tempfile
from datetime import datetime

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR

# ---------------------------------------------------------------------------
# Brand constants (mapped from GEN_COLORS)
# ---------------------------------------------------------------------------
BRAND = {
    "primary":   RGBColor(0x1B, 0x2A, 0x4A),  # deep navy
    "accent":    RGBColor(0xE6, 0x39, 0x46),  # signal red
    "success":   RGBColor(0x2E, 0xC4, 0xB6),  # teal
    "warning":   RGBColor(0xFF, 0x9F, 0x1C),  # amber
    "info":      RGBColor(0x45, 0x7B, 0x9D),  # steel blue
    "dark_text": RGBColor(0x1B, 0x2A, 0x4A),  # navy
    "muted":     RGBColor(0x6C, 0x75, 0x7D),  # gray
    "white":     RGBColor(0xFF, 0xFF, 0xFF),
    "light_bg":  RGBColor(0xF8, 0xF9, 0xFA),
    "grid":      RGBColor(0xE9, 0xEC, 0xEF),
}

RAG_COLORS = {
    "Green": RGBColor(0x2E, 0xC4, 0xB6),
    "Amber": RGBColor(0xFF, 0x9F, 0x1C),
    "Red":   RGBColor(0xE6, 0x39, 0x46),
}

# Slide dimensions (widescreen 16:9)
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Layout zones
TITLE_BAR_HEIGHT = Inches(1.1)
FOOTER_HEIGHT = Inches(0.45)
CONTENT_TOP = Inches(1.2)
CONTENT_LEFT = Inches(0.5)
CONTENT_WIDTH = Inches(12.333)
CONTENT_HEIGHT = SLIDE_HEIGHT - CONTENT_TOP - FOOTER_HEIGHT - Inches(0.2)

# Font sizes
TITLE_FONT_SIZE = Pt(28)
SUBTITLE_FONT_SIZE = Pt(14)
FOOTER_FONT_SIZE = Pt(10)
KPI_VALUE_SIZE = Pt(36)
KPI_LABEL_SIZE = Pt(12)
BODY_FONT_SIZE = Pt(14)
EXEC_SUMMARY_SIZE = Pt(16)
MIN_FONT_SIZE = Pt(14)


# ===========================================================================
# Low-level slide builders
# ===========================================================================

def _set_slide_size(prs):
    """Set presentation to widescreen 16:9."""
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT


def _add_title_bar(slide, title_text, subtitle_text=None):
    """Add a navy title bar to the top of a slide."""
    # Title bar background
    left, top = Inches(0), Inches(0)
    width = SLIDE_WIDTH
    height = TITLE_BAR_HEIGHT

    shape = slide.shapes.add_shape(
        1, left, top, width, height,  # 1 = rectangle
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = BRAND["primary"]
    shape.line.fill.background()

    # Title text
    txBox = slide.shapes.add_textbox(
        Inches(0.5), Inches(0.15), SLIDE_WIDTH - Inches(1), Inches(0.6),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = TITLE_FONT_SIZE
    p.font.bold = True
    p.font.color.rgb = BRAND["white"]
    p.alignment = PP_ALIGN.LEFT

    # Subtitle
    if subtitle_text:
        txBox2 = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.72), SLIDE_WIDTH - Inches(1), Inches(0.35),
        )
        tf2 = txBox2.text_frame
        tf2.word_wrap = True
        p2 = tf2.paragraphs[0]
        p2.text = subtitle_text
        p2.font.size = SUBTITLE_FONT_SIZE
        p2.font.italic = True
        p2.font.color.rgb = RGBColor(0xA0, 0xB0, 0xC0)
        p2.alignment = PP_ALIGN.LEFT


def _add_footer(slide, dataset_label, slide_num, total_slides):
    """Add footer with dataset label and page number."""
    footer_top = SLIDE_HEIGHT - FOOTER_HEIGHT

    # Left: dataset label
    txBox = slide.shapes.add_textbox(
        Inches(0.5), footer_top, Inches(8), FOOTER_HEIGHT,
    )
    tf = txBox.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = dataset_label
    p.font.size = FOOTER_FONT_SIZE
    p.font.color.rgb = BRAND["muted"]
    p.font.italic = True
    p.alignment = PP_ALIGN.LEFT

    # Right: page number
    txBox2 = slide.shapes.add_textbox(
        SLIDE_WIDTH - Inches(2), footer_top, Inches(1.5), FOOTER_HEIGHT,
    )
    tf2 = txBox2.text_frame
    tf2.vertical_anchor = MSO_ANCHOR.MIDDLE
    p2 = tf2.paragraphs[0]
    p2.text = f"{slide_num} / {total_slides}"
    p2.font.size = FOOTER_FONT_SIZE
    p2.font.color.rgb = BRAND["muted"]
    p2.alignment = PP_ALIGN.RIGHT


# ===========================================================================
# Slide type builders
# ===========================================================================

def add_title_slide(prs, storyline_config, cu_name, dataset_label, csm_name=""):
    """Create a branded title slide."""
    slide_layout = prs.slide_layouts[6]  # blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Full navy background
    bg = slide.shapes.add_shape(
        1, Inches(0), Inches(0), SLIDE_WIDTH, SLIDE_HEIGHT,
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = BRAND["primary"]
    bg.line.fill.background()

    # Accent line
    line = slide.shapes.add_shape(
        1, Inches(1.5), Inches(3.0), Inches(10.333), Inches(0.04),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = BRAND["success"]
    line.line.fill.background()

    # Title
    txBox = slide.shapes.add_textbox(
        Inches(1.5), Inches(1.5), Inches(10.333), Inches(1.3),
    )
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = storyline_config["title"]
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = BRAND["white"]
    p.alignment = PP_ALIGN.LEFT

    # Tagline
    txBox2 = slide.shapes.add_textbox(
        Inches(1.5), Inches(3.2), Inches(10.333), Inches(0.6),
    )
    tf2 = txBox2.text_frame
    tf2.word_wrap = True
    p2 = tf2.paragraphs[0]
    p2.text = storyline_config["tagline"]
    p2.font.size = Pt(22)
    p2.font.italic = True
    p2.font.color.rgb = BRAND["success"]
    p2.alignment = PP_ALIGN.LEFT

    # CU name
    txBox3 = slide.shapes.add_textbox(
        Inches(1.5), Inches(4.5), Inches(10.333), Inches(0.5),
    )
    tf3 = txBox3.text_frame
    p3 = tf3.paragraphs[0]
    p3.text = cu_name
    p3.font.size = Pt(20)
    p3.font.bold = True
    p3.font.color.rgb = BRAND["white"]
    p3.alignment = PP_ALIGN.LEFT

    # Dataset label + CSM
    txBox4 = slide.shapes.add_textbox(
        Inches(1.5), Inches(5.2), Inches(10.333), Inches(0.8),
    )
    tf4 = txBox4.text_frame
    tf4.word_wrap = True
    p4 = tf4.paragraphs[0]
    parts = [dataset_label]
    if csm_name:
        parts.append(f"Presented by {csm_name}")
    parts.append(datetime.now().strftime("%B %Y"))
    p4.text = " | ".join(parts)
    p4.font.size = Pt(14)
    p4.font.color.rgb = RGBColor(0xA0, 0xB0, 0xC0)
    p4.alignment = PP_ALIGN.LEFT

    # Audience
    txBox5 = slide.shapes.add_textbox(
        Inches(1.5), Inches(6.2), Inches(10.333), Inches(0.4),
    )
    tf5 = txBox5.text_frame
    p5 = tf5.paragraphs[0]
    p5.text = f"Audience: {storyline_config['audience']}"
    p5.font.size = Pt(11)
    p5.font.color.rgb = BRAND["muted"]
    p5.alignment = PP_ALIGN.LEFT

    return slide


def add_chart_slide(prs, title, subtitle, image_path,
                    dataset_label, slide_num, total_slides):
    """Add a slide with a chart image."""
    slide_layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(slide_layout)

    _add_title_bar(slide, title, subtitle)
    _add_footer(slide, dataset_label, slide_num, total_slides)

    # Insert chart image centered in content area
    if image_path and os.path.exists(image_path):
        # Calculate dimensions to fit while maintaining aspect ratio
        from PIL import Image as PILImage
        try:
            with PILImage.open(image_path) as img:
                img_w, img_h = img.size
        except Exception:
            img_w, img_h = 1600, 900

        aspect = img_w / img_h
        max_w = CONTENT_WIDTH
        max_h = CONTENT_HEIGHT

        if aspect > (max_w / max_h):
            # Width-constrained
            final_w = max_w
            final_h = int(max_w / aspect)
        else:
            # Height-constrained
            final_h = max_h
            final_w = int(max_h * aspect)

        # Center horizontally
        left = CONTENT_LEFT + (max_w - final_w) // 2
        top = CONTENT_TOP

        slide.shapes.add_picture(image_path, left, top, final_w, final_h)
    else:
        # Placeholder text when no image
        txBox = slide.shapes.add_textbox(
            CONTENT_LEFT, CONTENT_TOP + Inches(1.5),
            CONTENT_WIDTH, Inches(1),
        )
        tf = txBox.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = "[Chart will be generated when notebook data is available]"
        p.font.size = Pt(16)
        p.font.color.rgb = BRAND["muted"]
        p.font.italic = True
        p.alignment = PP_ALIGN.CENTER

    return slide


def add_exec_summary_slide(prs, storyline_config, kpi_values, findings,
                           dataset_label, slide_num, total_slides):
    """Add an executive summary slide with KPI cards and key findings."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_title_bar(slide, "Executive Summary",
                   storyline_config.get("tagline", ""))
    _add_footer(slide, dataset_label, slide_num, total_slides)

    # KPI cards row (up to 4 cards across the top)
    kpi_items = list(kpi_values.items())[:4]
    n_cards = len(kpi_items) if kpi_items else 0

    if n_cards > 0:
        card_width = Inches(2.8)
        card_height = Inches(1.6)
        card_spacing = Inches(0.3)
        total_card_width = (card_width * n_cards) + (card_spacing * (n_cards - 1))
        start_left = CONTENT_LEFT + (CONTENT_WIDTH - total_card_width) // 2

        for i, (label, value) in enumerate(kpi_items):
            left = start_left + i * (card_width + card_spacing)
            top = CONTENT_TOP + Inches(0.2)

            # Card background
            card = slide.shapes.add_shape(
                5, left, top, card_width, card_height,  # 5 = rounded rectangle
            )
            card.fill.solid()
            card.fill.fore_color.rgb = BRAND["light_bg"]
            card.line.color.rgb = BRAND["grid"]
            card.line.width = Pt(1)

            # Value
            vBox = slide.shapes.add_textbox(
                left + Inches(0.1), top + Inches(0.2),
                card_width - Inches(0.2), Inches(0.8),
            )
            vf = vBox.text_frame
            vf.word_wrap = True
            vp = vf.paragraphs[0]
            vp.text = str(value)
            vp.font.size = KPI_VALUE_SIZE
            vp.font.bold = True
            vp.font.color.rgb = BRAND["primary"]
            vp.alignment = PP_ALIGN.CENTER

            # Label
            lBox = slide.shapes.add_textbox(
                left + Inches(0.1), top + Inches(1.0),
                card_width - Inches(0.2), Inches(0.5),
            )
            lf = lBox.text_frame
            lf.word_wrap = True
            lp = lf.paragraphs[0]
            lp.text = label.replace("_", " ").title()
            lp.font.size = KPI_LABEL_SIZE
            lp.font.bold = True
            lp.font.color.rgb = BRAND["muted"]
            lp.alignment = PP_ALIGN.CENTER

    # Key findings section
    findings_top = CONTENT_TOP + Inches(2.2)
    if findings:
        # Section header
        hBox = slide.shapes.add_textbox(
            CONTENT_LEFT, findings_top, Inches(4), Inches(0.4),
        )
        hf = hBox.text_frame
        hp = hf.paragraphs[0]
        hp.text = "Key Findings"
        hp.font.size = Pt(18)
        hp.font.bold = True
        hp.font.color.rgb = BRAND["dark_text"]

        # Bullet points
        bBox = slide.shapes.add_textbox(
            CONTENT_LEFT + Inches(0.3), findings_top + Inches(0.5),
            Inches(11), Inches(3),
        )
        bf = bBox.text_frame
        bf.word_wrap = True

        for j, finding in enumerate(findings[:6]):
            p = bf.paragraphs[0] if j == 0 else bf.add_paragraph()
            p.text = finding
            p.font.size = EXEC_SUMMARY_SIZE
            p.font.color.rgb = BRAND["dark_text"]
            p.space_after = Pt(8)
            p.level = 0

    return slide


def add_text_slide(prs, title, subtitle, body_lines,
                   dataset_label, slide_num, total_slides):
    """Add a text-only slide (next steps, data sources, recommendations)."""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    _add_title_bar(slide, title, subtitle)
    _add_footer(slide, dataset_label, slide_num, total_slides)

    if body_lines:
        txBox = slide.shapes.add_textbox(
            CONTENT_LEFT + Inches(0.3), CONTENT_TOP + Inches(0.3),
            CONTENT_WIDTH - Inches(0.6), CONTENT_HEIGHT - Inches(0.3),
        )
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, line in enumerate(body_lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = line
            p.font.size = BODY_FONT_SIZE
            p.font.color.rgb = BRAND["dark_text"]
            p.space_after = Pt(10)

    return slide


# ===========================================================================
# Figure capture: re-execute cell code with plt.show() intercepted
# ===========================================================================

def render_cell_to_png(cell_ref, notebook_globals, output_dir, base_path):
    """Re-execute a cell file, intercept plt.show(), save figure to PNG.

    How it works:
    1. Read the cell source file from disk
    2. Monkey-patch plt.show() to savefig + plt.close() instead of display
    3. exec() the cell code using the notebook's live globals (so all
       DataFrames, palettes, and helpers are available)
    4. Restore plt.show()

    Args:
        cell_ref: e.g. "13-attrition/02_attrition_kpi"
        notebook_globals: the notebook's globals() dict (contains all
            DataFrames, GEN_COLORS, helper functions, etc.)
        output_dir: directory for output PNGs
        base_path: root path of the notebook cell files

    Returns:
        Path to saved PNG, or None on failure.
    """
    import matplotlib.pyplot as plt

    cell_path = os.path.join(base_path, cell_ref)
    if not os.path.isfile(cell_path):
        return None

    with open(cell_path, "r") as f:
        source = f.read()

    os.makedirs(output_dir, exist_ok=True)
    safe_name = cell_ref.replace("/", "_").replace("-", "_")
    png_path = os.path.join(output_dir, f"{safe_name}.png")

    # Track whether a figure was saved
    _captured = {"saved": False}

    # Monkey-patch plt.show to save instead of display
    _original_show = plt.show

    def _capture_show(*args, **kwargs):
        fig = plt.gcf()
        if fig and fig.get_axes():
            fig.savefig(png_path, dpi=300, bbox_inches="tight",
                        facecolor="white", edgecolor="none")
            _captured["saved"] = True
        plt.close("all")

    plt.show = _capture_show

    # Suppress display() calls (for HTML table cells)
    _original_display = notebook_globals.get("display")

    def _noop_display(*args, **kwargs):
        pass

    try:
        # Build exec namespace: notebook globals + our overrides
        exec_ns = dict(notebook_globals)
        exec_ns["display"] = _noop_display
        exec_ns["plt"] = plt

        exec(compile(source, cell_path, "exec"), exec_ns)

        # Copy any new variables back to notebook globals so downstream
        # cells that depend on computed values still work
        for k, v in exec_ns.items():
            if k not in ("__builtins__", "display", "plt"):
                notebook_globals[k] = v

    except Exception as e:
        # Cell failed -- not fatal, just skip this chart
        print(f"    WARNING: {cell_ref} failed: {e}")
        png_path = None
    finally:
        plt.show = _original_show
        if _original_display is not None:
            notebook_globals["display"] = _original_display
        plt.close("all")

    if _captured["saved"] and png_path and os.path.exists(png_path):
        return png_path
    return None


def resolve_chart_image(cell_ref, fig_registry, notebook_globals,
                        output_dir, base_path):
    """Get a chart PNG for a cell: use registry first, re-render if needed.

    Args:
        cell_ref: cell reference string
        fig_registry: dict of cell_ref -> PNG path (pre-captured)
        notebook_globals: notebook globals() for re-execution
        output_dir: temp dir for PNGs
        base_path: root path of cell files

    Returns:
        Path to PNG, or None.
    """
    # 1. Check if already in registry (pre-captured path)
    existing = fig_registry.get(cell_ref)
    if existing and os.path.exists(existing):
        return existing

    # 2. Re-render by executing the cell
    if notebook_globals and base_path:
        return render_cell_to_png(cell_ref, notebook_globals,
                                  output_dir, base_path)

    return None


# ===========================================================================
# Main export function
# ===========================================================================

def export_storyline(storyline_key, fig_registry=None, kpi_values=None,
                     findings=None, text_content=None,
                     cu_name="Credit Union", dataset_label="",
                     csm_name="", output_dir=None,
                     notebook_globals=None, base_path=None):
    """Export a single storyline deck to PPTX.

    Two modes of operation:
    1. Pre-captured: pass fig_registry with cell_ref -> PNG path mappings
    2. Re-render (recommended): pass notebook_globals=globals() and
       base_path pointing to the cell files root. Each chart cell is
       re-executed with plt.show() intercepted to save PNGs.

    Mode 2 is the default Jupyter workflow. After running all analysis
    cells, call this with notebook_globals=globals() and base_path set
    to the txn-visual-test root. The engine re-runs each chart cell in
    the notebook's namespace so all DataFrames are available.

    Args:
        storyline_key: key from STORYLINES dict
        fig_registry: dict of cell_ref -> PNG path (optional)
        kpi_values: dict of KPI label -> display value for exec summary
        findings: list of finding strings for exec summary
        text_content: dict of content_key -> list of body lines
        cu_name: credit union name
        dataset_label: date range string
        csm_name: CSM presenter name
        output_dir: directory for output PPTX
        notebook_globals: the notebook's globals() dict for re-rendering
        base_path: root directory of cell files (e.g., /tmp/txn-visual-test)

    Returns:
        Path to generated PPTX file.
    """
    from dashboard.storyline_config import STORYLINES

    config = STORYLINES.get(storyline_key)
    if config is None:
        raise ValueError(f"Unknown storyline: {storyline_key}")

    fig_registry = fig_registry or {}
    kpi_values = kpi_values or {}
    findings = findings or []
    text_content = text_content or {}
    output_dir = output_dir or os.getcwd()

    prs = Presentation()
    _set_slide_size(prs)

    slides = config["slides"]
    total_slides = len(slides)

    # Temp dir for chart PNGs
    tmp_dir = tempfile.mkdtemp(prefix="pptx_charts_")

    _rendered = 0
    _skipped = 0

    for slide_def in slides:
        slide_type = slide_def["type"]
        num = slide_def["num"]

        if slide_type == "title":
            add_title_slide(prs, config, cu_name, dataset_label, csm_name)

        elif slide_type == "exec_summary":
            add_exec_summary_slide(
                prs, config, kpi_values, findings,
                dataset_label, num, total_slides,
            )

        elif slide_type == "chart":
            cell_ref = slide_def.get("cell")
            img_path = None
            if cell_ref:
                img_path = resolve_chart_image(
                    cell_ref, fig_registry, notebook_globals,
                    tmp_dir, base_path,
                )
                if img_path:
                    _rendered += 1
                else:
                    _skipped += 1

            add_chart_slide(
                prs, slide_def["title"], slide_def.get("subtitle", ""),
                img_path, dataset_label, num, total_slides,
            )

        elif slide_type == "text":
            content_key = slide_def.get("content_key", "")
            body = text_content.get(content_key, [
                "Content to be customized by your CSM.",
                "Please add specific action items and timelines.",
            ])
            add_text_slide(
                prs, slide_def["title"], slide_def.get("subtitle", ""),
                body, dataset_label, num, total_slides,
            )

    # Save PPTX
    os.makedirs(output_dir, exist_ok=True)
    safe_key = storyline_key.replace(" ", "_")
    cu_safe = cu_name.replace(" ", "_").replace("/", "_")
    timestamp = datetime.now().strftime("%Y%m%d")
    filename = f"{config['id']:02d}_{safe_key}_{cu_safe}_{timestamp}.pptx"
    filepath = os.path.join(output_dir, filename)

    prs.save(filepath)
    print(f"    {config['title']}: {_rendered} charts rendered, "
          f"{_skipped} skipped, {total_slides} slides total")
    return filepath


def export_all_storylines(fig_registry=None, kpi_values=None,
                          findings_by_storyline=None,
                          text_content_by_storyline=None,
                          cu_name="Credit Union", dataset_label="",
                          csm_name="", output_dir=None,
                          core_only=False,
                          notebook_globals=None, base_path=None):
    """Export all (or core-only) storyline decks.

    Args:
        fig_registry: dict of cell_ref -> PNG path (optional)
        kpi_values: dict of KPI label -> display value (shared)
        findings_by_storyline: dict of storyline_key -> list of findings
        text_content_by_storyline: dict of storyline_key -> text_content dict
        cu_name: credit union name
        dataset_label: date range string
        csm_name: CSM name
        output_dir: output directory
        core_only: if True, skip SITUATIONAL storylines
        notebook_globals: notebook globals() for re-rendering charts
        base_path: root directory of cell files

    Returns:
        List of generated PPTX file paths.
    """
    from dashboard.storyline_config import PRESENTATION_ORDER, STORYLINES

    findings_by_storyline = findings_by_storyline or {}
    text_content_by_storyline = text_content_by_storyline or {}
    output_dir = output_dir or os.getcwd()

    paths = []
    for key in PRESENTATION_ORDER:
        config = STORYLINES[key]
        if core_only and config["type"] == "SITUATIONAL":
            continue

        findings = findings_by_storyline.get(key, [])
        text_content = text_content_by_storyline.get(key, {})

        path = export_storyline(
            key, fig_registry, kpi_values, findings, text_content,
            cu_name, dataset_label, csm_name, output_dir,
            notebook_globals=notebook_globals, base_path=base_path,
        )
        paths.append(path)

    return paths
